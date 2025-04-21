import streamlit as st
import pandas as pd
import plotly.express as px
import openai
import json
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
import datetime
import uuid
import gspread
import time
import requests
from requests.exceptions import RequestException
from google.oauth2.service_account import Credentials

# Validate OpenAI API key
if not st.secrets.get("OPENAI_API_KEY"):
    st.error("❌ OpenAI API key not configured.")
    st.stop()
openai.api_key = st.secrets["OPENAI_API_KEY"]

# --- Session State Initialization ---
if "chart_count" not in st.session_state:
    st.session_state.chart_count = 5
if "user_id" not in st.session_state:
    st.session_state.user_id = str(uuid.uuid4())

# --- Streamlit Page Config ---
st.set_page_config(page_title="Aarekha AI Charts", layout="wide")

# --- Set Full Background Color and Customize Styling ---
st.markdown(
    """
    <style>
        .stApp {
            background-color: #e0e7ff;
        }
        .stTextInput > div > input, .stTextArea > div > textarea {
            background-color: #ffffff !important;
            border: 2px solid #1e40af;
            border-radius: 8px;
            padding: 10px;
            color: #1e40af;
        }
        input[type="text"], textarea {
            background-color: #ffffff !important;
            color: #000000 !important;
            border: 2px solid #1A3C5A !important;
            border-radius: 8px !important;
            padding: 10px !important;
        }
        input::placeholder, textarea::placeholder {
            color: #000000 !important;
            opacity: 0.6 !important;
        }
        /* File Uploader Styling */
        .stFileUploader, 
        .stFileUploader label, 
        .stFileUploader [data-testid="stFileUploaderDropzone"], 
        .stFileUploader [data-testid="stFileUploaderDropzone"] p, 
        .stFileUploader [data-testid="stFileUploaderFileName"], 
        .stFileUploader [data-testid="stFileUploaderFileName"] p, 
        .stFileUploader [data-testid="stFileUploaderFileSize"], 
        .stFileUploader [data-testid="stFileUploaderFileSize"] p, 
        .stFileUploader [data-testid="stFileUploaderCloseButton"],
        .stFileUploader div, 
        .stFileUploader div div div, 
        .stFileUploader div[data-testid="stFileUploaderFile"],
        [data-testid="stFileUploader"] label,
        [data-testid="stFileUploader"] p {
            color: #000000 !important;
            background-color: white !important;
            border: 2px solid #1e40af;
            border-radius: 8px;
            padding: 15px;
            margin-bottom: 20px;
        }
        /* Color Picker Styling */
        .stColorPicker, 
        .stColorPicker label, 
        [data-testid="stColorPicker"] label,
        [data-testid="stColorPicker"] div > div > label {
            color: #000000 !important;
        }
        .stSidebar .stSlider label, 
        .stSidebar .stSlider [data-testid="stSliderValue"], 
        .stSidebar .stSlider span,
        .stSidebar .stSlider > div > div > div {
            color: #000000 !important;
        }
        h3, .stSelectbox label, 
        .stExpander label, 
        [data-testid="stExpander"] > div > div > div > label,
        [data-testid="stExpander"] > div > div > div > p,
        [data-testid="stExpander"] p, 
        [data-testid="stExpander"] div[role="button"] p {
            color: #000000 !important;
        }
        .stTextInput label, 
        .stTextArea label, 
        [data-testid="stForm"] > div > label {
            color: #000000 !important;
        }
        [data-testid="stSlider"] label, 
        [data-testid="stSlider"] span,
        [data-testid="stSliderValue"], 
        div[data-testid="stSlider"] div div div div,
        div[data-testid="stSlider"] span, 
        div[data-testid="stSlider"] div[data-testid="stTickBar"] div,
        div[data-testid="stSlider"] div[data-testid="stTickBar"] span,
        div[data-testid="stSlider"] div div span, 
        div[data-testid="stSlider"] .stTickBar span,
        div[data-testid="stSlider"] .stTickBar div, 
        div[data-testid="stSlider"] [data-testid="stTickMarks"] span,
        div[data-testid="stSlider"] [data-testid="stTickMarks"] div, 
        div[data-testid="stSlider"] * {
            color: #000000 !important;
        }
        [data-testid="stSpinner"] p {
            color: #000000 !important;
        }
        .stAlert, 
        [data-testid="stError"], 
        [data-testid="stWarning"], 
        [data-testid="alert-container"] * {
            color: #000000 !important;
        }
        [role="alert"] *, 
        .stException {
            color: #000000 !important;
        }
        .stButton > button {
            background-color: #1e40af;
            color: white;
            border-radius: 8px;
            padding: 10px 20px;
            border: none;
        }
        .stMarkdown h1, 
        .stMarkdown h2, 
        .stMarkdown p {
            color: #1e40af !important;
        }
        .element-container {
            margin: 10px 0;
        }
        .custom-card {
            background-color: white;
            padding: 20px;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
            margin-bottom: 20px;
        }
        .feedback-section {
            font-size: 24px;
            font-weight: bold;
            color: #1A3C5A;
        }
        .feedback-label {
            font-size: 14px;
            color: #6c6c6c;
        }
    </style>
    """,
    unsafe_allow_html=True
)

# --- Custom Title Section ---
st.markdown("""
    <div class='custom-card' style='text-align: center;'>
        <h1 style='font-size: 48px; color: #1e40af; margin: 0;'>AAREKHA</h1>
        <p style='font-size: 18px; color: #4b5563; margin: 10px 0 0;'>Your AI-Powered Chart Generator with Smart Insights</p>
    </div>
""", unsafe_allow_html=True)

# --- Early Access Alert ---
st.markdown("""
    <div class='custom-card' style='background: #ecfccb; padding: 15px; border-radius: 8px; text-align: center; color: #166534;'>
        🧪 <strong>Early Access:</strong> You're using the free launch version of <strong>Aarekha</strong>. Expect unlimited charts, PDF/PPT download, and more features soon!
    </div>
""", unsafe_allow_html=True)

# --- Upload Section ---
st.markdown("<h2 style='color: #1e40af; margin-bottom: 10px;'>📂 Upload Your Dataset</h2>", unsafe_allow_html=True)
st.markdown("<p style='color: #4b5563; margin-bottom: 15px;'>Upload a .csv or .xlsx file to get started. Aarekha will analyze your data and generate smart charts & insights for you.</p>", unsafe_allow_html=True)

file = st.file_uploader("Upload CSV or Excel file", type=["csv", "xlsx"])

# --- Helper Functions ---
def auto_group_large_data(df, x_col):
    if df[x_col].nunique() > 50:
        return df.groupby(x_col).size().sort_values(ascending=False).head(20).index.tolist()
    return df[x_col].unique()

def save_feedback_to_gsheet(email, feedback):
    try:
        if "gspread" not in st.secrets:
            raise KeyError("st.secrets has no key 'gspread'. Please configure Google service account credentials in Streamlit Cloud secrets.")
        scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
        creds = Credentials.from_service_account_info(st.secrets["gspread"], scopes=scope)
        client = gspread.authorize(creds)
        sheet = client.open("Aarekha Feedback").sheet1
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        chart_count = st.session_state.get("chart_count", "N/A")
        sheet.append_row([timestamp, email, chart_count, feedback])
    except Exception as e:
        st.error(f"❌ Failed to save feedback: {str(e)}")

# --- Main Logic ---
if file:
    try:
        df = pd.read_csv(file, encoding='latin1') if file.name.endswith(".csv") else pd.read_excel(file)
    except Exception as e:
        st.error(f"❌ Failed to load file: {e}")
        st.stop()

    # Preprocess: Create a count-based column for datasets without numeric columns
    numeric_columns = [col for col in df.columns if pd.api.types.is_numeric_dtype(df[col])]
    if not numeric_columns:
        df['Order Count'] = 1
        numeric_columns = ['Order Count']
    
    st.dataframe(df.head(), use_container_width=True)

    st.sidebar.header("🔍 Global Filters")
    filter_columns = ["Date", "Region", "Category"]
    filter_values = {}

    for col in filter_columns:
        if col in df.columns:
            unique_vals = df[col].dropna().unique()
            selection = st.sidebar.multiselect(f"Filter by {col}", options=unique_vals, default=unique_vals, key=f"global_filter_{col}")
            filter_values[col] = selection

    filtered_df = df.copy()
    for col, vals in filter_values.items():
        filtered_df = filtered_df[filtered_df[col].isin(vals)]

    num_charts = st.slider("How many charts to auto-generate?", 1, 10, 5)

    if "chart_plans" not in st.session_state or "data_hash" not in st.session_state or st.session_state.data_hash != hash(df.head(100).to_csv(index=False)) or st.session_state.chart_count != num_charts:
        system_msg = """
You are a senior data analyst helping businesses understand their data through effective visual storytelling.
Based on the sample dataset provided, intelligently recommend the most appropriate charts and write a meaningful, user-friendly insight for each.
For each chart:
- Choose the most suitable chart type from this list: Bar, Line, Scatter, Pie, Histogram, Heatmap, Box, Area, Bubble, Stacked Bar.
- Pick logical x-axis and y-axis columns based on data patterns (e.g., time series, categories, quantities, numerical values).
- If the dataset lacks numeric columns, use count-based aggregations (e.g., count of orders by category, date, or other categorical columns) for y-axis or values, and prioritize Bar, Pie, Histogram, or Heatmap charts.
- For Line, Scatter, Box, Area, and Bubble charts, ensure the y-axis is a numeric column (e.g., sales, revenue, quantity, or aggregated counts). If no numeric columns exist, avoid these chart types.
- Write the insight as **three clearly separated bullet points** in a brief tone:
  - **Key Observation:** A brief summary of what the chart reveals (e.g., trend, peak, dip, comparison, correlation).
  - **Business Impact:** What this trend or insight means for the business. Focus on how it helps or hurts growth, revenue, efficiency, etc.
  - **Recommended Action:** Suggest a clear, practical step the business can take based on the insight (e.g., optimize, investigate, invest, improve, explore further).
Respond only in this valid JSON list format, ensuring proper escaping of special characters (e.g., newlines, tabs):
[
  {
    "chart_type": "Bar",
    "x": "column_name",
    "y": "Order Count",
    "insight": "- **Key Observation:** Description.\\n- **Business Impact:** Description.\\n- **Recommended Action:** Description."
  }
]
"""
        sample = df.head(100).to_csv(index=False)
        try:
            with st.spinner("🤖 Analyzing data & generating chart recommendations..."):
                max_retries = 3
                for attempt in range(max_retries):
                    try:
                        res = openai.ChatCompletion.create(
                            model="gpt-3.5-turbo",
                            messages=[
                                {"role": "system", "content": system_msg},
                                {"role": "user", "content": f"Here is a data sample:\n{sample}\n\nGenerate {num_charts} intelligent chart recommendations."}
                            ],
                            temperature=0.4,
                            max_tokens=1500,
                            timeout=30
                        )
                        # Clean the response to remove code block markers and control characters
                        raw_response = res.choices[0].message.content.strip()
                        if raw_response.startswith("```"):
                            raw_response = raw_response.split("\n", 1)[1].rsplit("\n", 1)[0]
                        raw_response = raw_response.replace("\r", "").replace("\t", " ")
                        st.session_state.chart_plans = json.loads(raw_response)
                        st.session_state.insights = [plan.get("insight", "") for plan in st.session_state.chart_plans]
                        st.session_state.data_hash = hash(sample)
                        st.session_state.chart_count = num_charts
                        break
                    except json.JSONDecodeError as je:
                        if attempt < max_retries - 1:
                            st.warning(f"⚠️ Attempt {attempt + 1} failed due to invalid JSON response: {je}. Retrying in 5 seconds...")
                            time.sleep(5)
                        else:
                            st.error(f"❌ AI chart recommendation failed after {max_retries} attempts due to invalid JSON response: {je}. Falling back to default charts.")
                            st.session_state.chart_plans = [
                                {
                                    "chart_type": "Bar",
                                    "x": df.columns[0],
                                    "y": "Order Count",
                                    "insight": "- **Key Observation:** Default chart generated.\n- **Business Impact:** Limited analysis due to API error.\n- **Recommended Action:** Try again later."
                                }
                            ] * num_charts
                            st.session_state.insights = [plan["insight"] for plan in st.session_state.chart_plans]
                    except requests.exceptions.RequestException as re:
                        if attempt < max_retries - 1:
                            st.warning(f"⚠️ Attempt {attempt + 1} failed due to network issue: {re}. Retrying in 5 seconds...")
                            time.sleep(5)
                        else:
                            st.error(f"❌ AI chart recommendation failed after {max_retries} attempts: {re}. Check your internet connection, API key, or OpenAI status.")
                            st.session_state.chart_plans = [
                                {
                                    "chart_type": "Bar",
                                    "x": df.columns[0],
                                    "y": "Order Count",
                                    "insight": "- **Key Observation:** Default chart generated.\n- **Business Impact:** Limited analysis due to API error.\n- **Recommended Action:** Try again later."
                                }
                            ] * num_charts
                            st.session_state.insights = [plan["insight"] for plan in st.session_state.chart_plans]
        except Exception as e:
            st.error(f"❌ Unexpected error during chart generation: {e}")
            st.session_state.chart_plans = [
                {
                    "chart_type": "Bar",
                    "x": df.columns[0],
                    "y": "Order Count",
                    "insight": "- **Key Observation:** Default chart generated.\n- **Business Impact:** Limited analysis due to error.\n- **Recommended Action:** Try again later."
                }
            ] * num_charts
            st.session_state.insights = [plan["insight"] for plan in st.session_state.chart_plans]

    chart_plans = st.session_state.chart_plans
    insights = st.session_state.insights

    chart_type_options = ["Bar", "Line", "Scatter", "Pie", "Histogram", "Heatmap", "Box", "Area", "Bubble", "Stacked Bar"]
    color_palette = px.colors.qualitative.Set3

    ppt = Presentation()
    blank_slide_layout = ppt.slide_layouts[6]
    pdf_buffer = io.BytesIO()
    c = canvas.Canvas(pdf_buffer, pagesize=letter)

    for idx, plan in enumerate(chart_plans):
        with st.container():
            st.markdown(f"### Chart {idx+1}")
            col1, col2, col3 = st.columns(3)

            default_chart = plan['chart_type'] if plan['chart_type'] in chart_type_options else "Bar"
            default_x = plan.get('x') if plan.get('x') in df.columns else df.columns[0]
            default_y = plan.get('y') if plan.get('y') in df.columns else numeric_columns[0]

            chart_type = col1.selectbox(f"Chart Type {idx+1}", chart_type_options, index=chart_type_options.index(default_chart), key=f"type_{idx}")
            x_axis = col2.selectbox(f"X-axis {idx+1}", df.columns, index=df.columns.get_loc(default_x), key=f"x_{idx}")
            
            # Deduplicate y_axis_options to prevent multiple 'Order Count'
            if chart_type in ["Line", "Scatter", "Box", "Area", "Bubble"]:
                y_axis_options = numeric_columns
            else:
                y_axis_options = list(dict.fromkeys(df.columns.tolist()))  # Remove duplicates while preserving order
            y_axis = col3.selectbox(
                f"Y-axis {idx+1}",
                y_axis_options,
                index=y_axis_options.index(default_y) if default_y in y_axis_options else 0,
                key=f"y_{idx}"
            ) if chart_type not in ["Pie", "Histogram"] else None
            
            color = col1.color_picker(f"Pick a chart color for Chart {idx+1}", value="#1f77b4", key=f"color_{idx}")

            with st.expander(f"🔧 Optional Filters for Chart {idx+1}"):
                perchart_filters = {}
                for col in df.columns:
                    if df[col].nunique() < 20 and df[col].dtype == "object":
                        selected = st.multiselect(f"Filter {col}", df[col].unique(), default=df[col].unique(), key=f"filter_{col}_{idx}")
                        perchart_filters[col] = selected

                chart_df = filtered_df.copy()
                for col, vals in perchart_filters.items():
                    chart_df = chart_df[chart_df[col].isin(vals)]

                if chart_df.empty:
                    st.warning(f"⚠️ Chart {idx+1} has no data after filtering. Please adjust filters.")
                    continue

            fig = None
            try:
                if chart_type == "Bar":
                    if y_axis and y_axis in numeric_columns:
                        fig = px.bar(chart_df, x=x_axis, y=y_axis, color_discrete_sequence=[color], hover_data=[x_axis, y_axis])
                    else:
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Histogram":
                    fig = px.histogram(chart_df, x=x_axis, color_discrete_sequence=[color], hover_data=[x_axis])
                elif chart_type == "Scatter":
                    if y_axis and y_axis in numeric_columns:
                        fig = px.scatter(chart_df, x=x_axis, y=y_axis, color_discrete_sequence=[color], hover_data=[x_axis, y_axis])
                    else:
                        st.warning(f"⚠️ Chart {idx+1} failed: Y-axis '{y_axis}' is not numeric. Using bar chart instead.")
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Line":
                    if y_axis and y_axis in numeric_columns:
                        if chart_df[x_axis].nunique() > 20:
                            chart_df = chart_df.groupby(x_axis).agg({y_axis: 'mean'}).reset_index()
                        chart_df = chart_df.nlargest(10, y_axis)
                        fig = px.line(chart_df, x=x_axis, y=y_axis, color_discrete_sequence=[color], hover_data=[x_axis, y_axis])
                    else:
                        st.warning(f"⚠️ Chart {idx+1} failed: Y-axis '{y_axis}' is not numeric. Using bar chart instead.")
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Box":
                    if y_axis and y_axis in numeric_columns:
                        fig = px.box(chart_df, x=x_axis, y=y_axis, color_discrete_sequence=[color], hover_data=[x_axis, y_axis])
                    else:
                        st.warning(f"⚠️ Chart {idx+1} failed: Y-axis '{y_axis}' is not numeric. Using bar chart instead.")
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Pie":
                    pie_data = chart_df[x_axis].value_counts().reset_index()
                    pie_data.columns = [x_axis, 'count']
                    fig = px.pie(pie_data, names=x_axis, values='count', color_discrete_sequence=px.colors.qualitative.Set3, hover_data=[x_axis, 'count'])
                elif chart_type == "Area":
                    if y_axis and y_axis in numeric_columns:
                        fig = px.area(chart_df, x=x_axis, y=y_axis, color_discrete_sequence=[color], hover_data=[x_axis, y_axis])
                    else:
                        st.warning(f"⚠️ Chart {idx+1} failed: Y-axis '{y_axis}' is not numeric. Using bar chart instead.")
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Heatmap":
                    if y_axis:
                        pivot = chart_df.pivot_table(index=x_axis, columns=y_axis, aggfunc='size', fill_value=0)
                        fig = px.imshow(pivot, color_continuous_scale='Viridis', hover_data=None)
                    else:
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Bubble":
                    if y_axis and y_axis in numeric_columns:
                        fig = px.scatter(chart_df, x=x_axis, y=y_axis, size=y_axis, color_discrete_sequence=[color], hover_data=[x_axis, y_axis])
                    else:
                        st.warning(f"⚠️ Chart {idx+1} failed: Y-axis '{y_axis}' is not numeric. Using bar chart instead.")
                        count_df = chart_df.groupby(x_axis).size().reset_index(name='Order Count')
                        fig = px.bar(count_df, x=x_axis, y='Order Count', color_discrete_sequence=[color], hover_data=[x_axis, 'Order Count'])
                elif chart_type == "Stacked Bar":
                    if y_axis and y_axis in numeric_columns:
                        fig = px.bar(chart_df, x=x_axis, y=y_axis, color=x_axis, barmode='stack', color_discrete_sequence=px.colors.qualitative.Pastel, hover_data=[x_axis, y_axis])
                    else:
                        count_df = chart_df.groupby([x_axis, y_axis]).size().unstack(fill_value=0).reset_index()
                        fig = px.bar(count_df, x=x_axis, y=count_df.columns[1:], barmode='stack', color_discrete_sequence=px.colors.qualitative.Pastel, hover_data=[x_axis])

                if fig:
                    st.plotly_chart(fig, use_container_width=True, key=f"plotly_chart_{idx}_{uuid.uuid4()}")

                with st.expander(f"🔎 Show Data for Chart {idx+1}"):
                    display_cols = [x_axis] + ([y_axis] if y_axis and y_axis in chart_df.columns else [])
                    st.dataframe(chart_df[display_cols].dropna().head(50), use_container_width=True)

                regenerate = st.button(f"🔁 Regenerate Insight for Chart {idx+1}", key=f"regen_{idx}")

                if regenerate:
                    try:
                        prompt = f"""
                        Given this chart config with x={x_axis}, y={y_axis or 'count'}, type={chart_type}, generate a business insight using the 3-bullet format in a brief tone.
                        Data sample:\n{chart_df.head(100).to_csv(index=False)}
                        Respond in valid JSON format, ensuring proper escaping of special characters:
                        {
                          "insight": "- **Key Observation:** Description.\\n- **Business Impact:** Description.\\n- **Recommended Action:** Description."
                        }
                        """
                        result = openai.ChatCompletion.create(
                            model="gpt-3.5-turbo",
                            messages=[
                                {"role": "system", "content": system_msg},
                                {"role": "user", "content": prompt}
                            ],
                            temperature=0.3,
                            max_tokens=600,
                            timeout=30
                        )
                        # Clean the insight response
                        raw_insight = result.choices[0].message.content.strip()
                        if raw_insight.startswith("```"):
                            raw_insight = raw_insight.split("\n", 1)[1].rsplit("\n", 1)[0]
                        raw_insight = raw_insight.replace("\r", "").replace("\t", " ")
                        insight_data = json.loads(raw_insight)
                        insights[idx] = insight_data["insight"]
                        st.session_state.insights = insights
                    except json.JSONDecodeError as je:
                        st.warning(f"⚠️ Insight regeneration for Chart {idx+1} failed due to invalid JSON: {je}. Keeping previous insight.")
                    except requests.exceptions.RequestException as re:
                        st.warning(f"⚠️ Insight regeneration for Chart {idx+1} failed due to network issue: {re}. Keeping previous insight.")
                    except Exception as e:
                        st.warning(f"⚠️ Insight regeneration for Chart {idx+1} failed: {e}. Keeping previous insight.")

                st.markdown(f"<div class='custom-card' style='padding: 10px; font-size: 14px; color: #333'><strong>Insight:</strong><br>{insights[idx].replace(chr(10), '<br>')}</div>", unsafe_allow_html=True)

                if fig:
                    img_bytes = fig.to_image(format="png")
                    st.download_button(
                        label="Download PNG",
                        data=img_bytes,
                        file_name=f"chart_{idx+1}.png",
                        mime="image/png",
                        key=f"dl_chart_{idx}"
                    )

                    slide = ppt.slides.add_slide(blank_slide_layout)
                    image_stream = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(image_stream, Inches(1), Inches(0.5), height=Inches(4.5))
                    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(8.5), Inches(2.5))
                    tf = textbox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = insights[idx].replace("\n", "\n")
                    font = run.font
                    font.size = Pt(12)
                    font.name = 'Arial'
                    font.color.rgb = RGBColor(50, 50, 50)

                    image_stream.seek(0)
                    img_reader = ImageReader(image_stream)
                    c.drawImage(img_reader, 50, 400, width=500, height=250)
                    c.setFont("Helvetica", 10)
                    c.drawString(50, 380, insights[idx].replace("\n", " ")[:400])
                    c.showPage()

            except Exception as e:
                st.warning(f"⚠️ Chart {idx+1} failed: {e}")

    ppt_buffer = io.BytesIO()
    ppt.save(ppt_buffer)
    st.download_button("📊 Download Full Report as PPT", data=ppt_buffer.getvalue(), file_name="Aarekha_Charts_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    c.save()
    pdf_buffer.seek(0)
    st.download_button("📄 Download Full Report as PDF", data=pdf_buffer, file_name="Aarekha_Charts_Report.pdf", mime="application/pdf")

# --- Unified Email + Feedback Section ---
st.markdown("<h2 style='color: #1e40af; margin-bottom: 10px;'>📬 Please Share Feedback</h2>", unsafe_allow_html=True)
with st.form("email_feedback_form"):
    email = st.text_input("📧 Enter your email (required to use the app):", key="user_email")
    feedback = st.text_area("💬 We'd love your feedback! What do you think, what could be better, or what feature you need most:")
    submit_both = st.form_submit_button("Submit")
    if submit_both:
        if email:
            save_feedback_to_gsheet(email, feedback)
            st.success("✅ Thanks for your feedback!")
        else:
            st.warning("⚠️ Please enter your email to continue.")
